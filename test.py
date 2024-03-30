import aspose.slides as slides
import aspose.pydrawing as drawing
import aspose.words as words
import pptx

# with slides.Presentation("./subhadipjana866@gmail.com/network/ATM.ppt") as presentation:
#     presentation.save("presentation.pdf", slides.export.SaveFormat.PDF)

ppt = pptx.Presentation("./subhadipjana866@gmail.com/network/ATM.ppt")
text = []
for slide in ppt.slides:
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            text.append(shape.text)
with open("./subhadipjana866@gmail.com/network/ATM.txt", "w") as f:
    f.write("\n".join(text))

# presentation = slides.Presentation("./subhadipjana866@gmail.com/network/MODEM.pptx")
# doc = words.Document()
# builder = words.DocumentBuilder(doc)

# for index in range(presentation.slides.length):
#     slide = presentation.slides[index]
    
#     for shape in slide.shapes:
#         # inserts slide's texts
#         if (type(shape) is slides.AutoShape):
#             builder.writeln(shape.text_frame.text)
   
#     builder.insert_break(words.BreakType.PAGE_BREAK)
# doc.save("presentation.docx")